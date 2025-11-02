import os
from pathlib import Path

import openai
from fastapi import FastAPI, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.util import Inches
from backend.format_ppt import apply_formatting
from backend.requirement_enricher import RequirementEnricher
from backend.db_handler import store_ai_feedback, retrieve_common_feedback

# ------------------------- 🚀 Initialize FastAPI App -------------------------
app = FastAPI()

@app.get("/")
def home():
    return {"message": "AI Presentation Generator is Running!"}

# ------------------------- 📁 File Paths -------------------------
BASE_DIR = Path(__file__).resolve().parents[1]
OUTPUT_DIR = BASE_DIR / "output"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# ------------------------- 🤖 Load OpenAI API Key -------------------------
openai_api_key = os.getenv("OPENAI_API_KEY")
if not openai_api_key:
    raise ValueError("OpenAI API key not found. Set OPENAI_API_KEY as an environment variable.")

client = openai.OpenAI(api_key=openai_api_key)
enricher = RequirementEnricher(api_key=openai_api_key)

SYSTEM_PROMPT = (
    "You are an expert presentation writer who creates concise, bulleted slide content. "
    "Respond only with the body copy for the slide. Do not include slide numbers, titles, "
    "or any assistant preamble or postscript."
)

# ------------------------- 📄 Request Model -------------------------
class PresentationRequest(BaseModel):
    topic: str = Field(..., example="AI in Finance")
    num_slides: int = Field(..., ge=1, le=20, example=5)
    audience: str = Field(default="General Public")
    duration: int = Field(default=20)
    purpose: str = Field(default="Explain how AI is used in finance.")
    design_style: str = Field(default="Minimalist")
    font_choice: str = Field(default="Arial")
    color_scheme: str = Field(default="#000000")  # Used as font color now
    additional_notes: str = Field(default="")

# ------------------------- 📝 Generate PPT with AI Optimization -------------------------
@app.post("/generate_ppt")
def generate_ppt(request: PresentationRequest):
    try:
        print(f"🟢 Generating PPT for topic: {request.topic} | Slides: {request.num_slides}")

        filename = f"{request.topic.replace(' ', '_')}_presentation.pptx"
        file_path = OUTPUT_DIR / filename

        if file_path.exists():
            file_path.unlink()

        prs = Presentation()

        # ✅ AI-Generated Slide Titles (Enforcing Slide Count)
        enriched_titles = enricher.generate_slide_titles(request.topic, request.num_slides)

        # ✅ AI-Optimized Prompt for Content (Enforcing Slide Count)
        refined_prompt = enricher.enrich_prompt(
            request.topic, request.audience, request.duration, request.purpose, request.num_slides
        )

        # ✅ GPT-Generated Slide Content (Batch Processing + Slide Count Fix)
        slides_prompts = [
            {
                "role": "user",
                "content": f"Slide {i+1}: {enriched_titles[i]}\n{refined_prompt}\nEnsure {request.num_slides} slides."
            }
            for i in range(request.num_slides)
        ]

        slide_contents = []
        for idx, user_prompt in enumerate(slides_prompts):
            try:
                response = client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": SYSTEM_PROMPT},
                        user_prompt,
                    ],
                )
                content = response.choices[0].message.content.strip()
                if not content:
                    raise ValueError(f"Empty response for slide {idx + 1}.")
                slide_contents.append(content)
            except Exception as api_error:
                raise HTTPException(status_code=500, detail=f"❌ GPT API Error: {str(api_error)}") from api_error

        # ✅ Ensure Slide Content Matches Requested Count
        if len(slide_contents) != request.num_slides:
            raise HTTPException(
                status_code=500,
                detail=f"⚠️ AI returned {len(slide_contents)} slides instead of {request.num_slides}. Please retry."
            )

        # ✅ Generate Slides with AI-Formatted Content
        for i, slide_content in enumerate(slide_contents):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = enriched_titles[i]
            else:
                title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
                title_box.text_frame.text = enriched_titles[i]

            left, top, width, height = Inches(1), Inches(1.5), Inches(8), Inches(4)
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.text = slide_content

            text_frame.margin_left = Inches(0.2)
            text_frame.margin_right = Inches(0.2)
            text_frame.margin_top = Inches(0.2)
            text_frame.margin_bottom = Inches(0.2)
            text_frame.word_wrap = True

            # ✅ Store AI Feedback for Continuous Improvement
            store_ai_feedback(request.topic, i+1, slide_content)

        # ✅ Apply AI-Driven Formatting
        user_preferences = {
            "font_choice": request.font_choice,
            "color_scheme": request.color_scheme,
        }
        apply_formatting(prs, user_preferences)

        prs.save(str(file_path))
        print(f"✅ Presentation saved successfully: {file_path}")

        return {"message": "✅ Presentation created successfully", "file": filename}

    except Exception as e:
        print(f"❌ Error generating presentation: {str(e)}")
        raise HTTPException(status_code=500, detail=f"❌ Error generating presentation: {str(e)}")


# ------------------------- 📥 Smart Backend Preview -------------------------
@app.get("/preview_ppt/{filename}")
def preview_ppt(filename: str):
    """Returns the text content of the generated PPT for quick review & improvement."""
    file_path = OUTPUT_DIR / filename
    
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="❌ File not found")

    ppt = Presentation(str(file_path))
    preview = []
    for i, slide in enumerate(ppt.slides):
        text = "\n".join(shape.text for shape in slide.shapes if shape.has_text_frame)
        preview.append(f"Slide {i+1}:\n{text}")

    return {"preview": "\n\n".join(preview)}


# ------------------------- 📥 Download PPT -------------------------
@app.get("/download_ppt/{filename}")
def download_ppt(filename: str):
    file_path = OUTPUT_DIR / filename
    if not file_path.exists():
        raise HTTPException(status_code=404, detail=f"❌ File '{filename}' not found.")
    
    return FileResponse(str(file_path), media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename=filename)

# ------------------------- 🏁 Start API -------------------------
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000, reload=True)
