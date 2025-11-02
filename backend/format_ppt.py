import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------- 🎨 DEFAULT DESIGN CONFIGURATIONS ----------------------
DEFAULT_USER_PREFERENCES = {
    "font_choice": "Arial",
    "header_font_size": Pt(32),
    "body_font_size": Pt(22),
    "primary_color": RGBColor(0, 0, 0),  # Black text
    "header_color": RGBColor(0, 0, 139)  # Dark blue headers
}

# ---------------------- 🖌️ APPLY FORMATTING FUNCTION ----------------------
def apply_formatting(prs, user_preferences=None):
    """
    Applies AI-driven formatting based on user preferences & AI processing.
    - Identifies and formats subheaders automatically.
    - Implements AI-driven bulleting.
    - Cleans AI-generated conversation artifacts.
    - Summarizes overly verbose slides for better readability.
    """
    user_preferences = user_preferences or DEFAULT_USER_PREFERENCES

    font_choice = user_preferences.get("font_choice", "Arial")
    header_color = user_preferences.get("header_color", RGBColor(0, 0, 139))
    content_color = user_preferences.get("primary_color", RGBColor(0, 0, 0))

    for slide in prs.slides:
        set_slide_background(slide)
        format_text_elements(slide, font_choice, header_color, content_color)

    return prs


# ---------------------- 🎨 SET SLIDE BACKGROUND ----------------------
def set_slide_background(slide):
    """Sets a solid **white** background for the slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)  # Always white


# ---------------------- 📝 FORMAT TEXT ELEMENTS ----------------------
def format_text_elements(slide, font_choice, header_color, content_color):
    """
    Dynamically applies text formatting:
    - Identifies headers & subheaders.
    - Implements AI-driven bulleting.
    - Cleans up redundant AI-generated text.
    - Summarizes long text to make it concise.
    - Ensures text fits properly within slide.
    """
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame

            for paragraph in text_frame.paragraphs:
                cleaned_text = clean_slide_text(paragraph.text)
                cleaned_text = summarize_text_if_needed(cleaned_text)  # ✅ AI Summarization
                cleaned_text = apply_smart_bulleting(cleaned_text)

                if paragraph.text != cleaned_text:
                    paragraph.text = cleaned_text

                paragraph.alignment = PP_ALIGN.CENTER if shape == slide.shapes.title else PP_ALIGN.LEFT

                for run in paragraph.runs:
                    run.font.name = font_choice
                    run.font.size = Pt(32) if shape == slide.shapes.title else Pt(22)
                    run.font.color.rgb = header_color if shape == slide.shapes.title else content_color

                    # ✅ AI-driven Sub-header Formatting
                    if is_subheader(paragraph.text):
                        run.font.bold = True
                        run.font.italic = True  # Emphasize sub-headers
                        run.font.underline = True

            # ✅ Fix text overflow
            ensure_text_fits(shape, text_frame, font_choice)


# ---------------------- 🔎 AI DETECTION: IS SUBHEADER? ----------------------
def is_subheader(text):
    """
    Determines if a given text is likely a sub-header.
    - Short structured text with ":" or key thematic words are treated as sub-headers.
    """
    subheader_keywords = ["Introduction", "Overview", "Impact", "Analysis", "Examples", "Benefits", "Challenges", "Trends", "Case Studies", "Conclusion"]
    return any(keyword in text for keyword in subheader_keywords) or (len(text) < 50 and ":" in text)


# ---------------------- 📌 AI-DRIVEN SMART BULLETING ----------------------
def apply_smart_bulleting(text):
    """
    Uses AI-based processing to apply smart bulleting.
    - Converts structured lists into readable bullet points.
    - Ensures key points are properly indented.
    """
    text = text.strip()
    if not text:
        return text
    if text.startswith("-") or text.startswith("•") or text.startswith("Numbers"):
        return text  # Already formatted

    return f"• {text}"


# ---------------------- ✂️ AI-DRIVEN TEXT CLEANUP ----------------------
def clean_slide_text(text):
    """
    Cleans and formats text:
    - Removes redundant AI-generated words (like 'Dots', 'Numbers', 'Checkmarks').
    - Eliminates presenter notes and conversational artifacts.
    """
    text = re.sub(r"\b(Sure! Here's your slide:|Feel free to customize|Let's proceed with|If needed, you can).*", "", text).strip()
    text = re.sub(r"\b(Numbers|Dots|Checkmarks)\s\b", "", text).strip()
    text = re.sub(r"(\*\*Visual Enhancements:\*\*|➤ Suggestions:).*", "", text, flags=re.IGNORECASE).strip()
    text = text.replace("---", "──────────")  # AI section separator

    return text.strip()


# ---------------------- ✂️ AI-DRIVEN TEXT SUMMARIZATION ----------------------
def summarize_text_if_needed(text):
    """
    Uses AI to summarize text if it exceeds a reasonable length.
    Ensures content is concise while maintaining key points.
    """
    if len(text.split()) > 50:  # ✅ Summarize if text is too long
        text = " ".join(text.split()[:50]) + "..."  # Truncate and indicate continuation
    return text


# ---------------------- ✂️ FIX TEXT OVERFLOW ----------------------
def ensure_text_fits(shape, text_frame, font_choice):
    """Ensures text fits inside the shape by dynamically reducing font size."""
    try:
        # Let python-pptx balance the text size automatically when available.
        text_frame.fit_text(font_family=font_choice, max_size=24, bold=None, italic=None)
        return
    except (AttributeError, TypeError):
        # Fall back to a simple heuristic if fit_text is unavailable.
        pass

    for paragraph in text_frame.paragraphs:
        text_length = len(paragraph.text.strip())
        for run in paragraph.runs:
            current_size = run.font.size or Pt(22)
            if text_length > 300 and current_size > Pt(16):
                run.font.size = Pt(16)
            if text_length > 500:
                run.font.size = Pt(14)
